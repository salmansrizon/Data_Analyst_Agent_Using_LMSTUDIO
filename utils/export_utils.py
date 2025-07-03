import pandas as pd
from pptx import Presentation
from pptx.util import Inches


def save_to_csv(df, output_path):
    df.to_csv(output_path, index=False)


def create_pptx(analysis, visuals, output_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Data Analysis Report"

    for visual in visuals:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        pic = slide.shapes.add_picture(visual, Inches(1), Inches(1), width=Inches(6))

    prs.save(output_path)
